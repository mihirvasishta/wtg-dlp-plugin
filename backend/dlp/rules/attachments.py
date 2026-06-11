"""
Rule 3 — Attachment Inspection

Six-tier pipeline (all in-memory, no disk writes):

  Tier 1 — Extension block  : hard-blocked extensions → severity BLOCK
  Tier 2 — Size limits      : warn / block on file size
  Tier 3 — Plain-text scan  : decode text/csv/etc. → regex patterns
  Tier 4 — PDF extraction   : pypdf → text → regex patterns
  Tier 5 — Office extraction: python-docx / openpyxl / python-pptx → text
  Tier 6 — ZIP inspection   : member names; optional content extraction

Returns:
  - A list of Violation objects
  - A list of extracted text strings (one per attachment, empty string if
    extraction failed or not applicable) — used by Rule 2 for partner-name
    scanning inside attachment content.
"""
from __future__ import annotations
import base64
import io
import os
import zipfile
from typing import List, Tuple

from models import AttachmentPayload, Violation

# Default attachment policy (overridable in future config)
BLOCKED_EXTENSIONS = {".exe", ".bat", ".cmd", ".ps1", ".vbs", ".scr", ".msi", ".dll"}
WARN_EXTENSIONS = {".zip", ".rar", ".7z", ".iso"}
WARN_SIZE_MB = 10
BLOCK_SIZE_MB = 25
SCAN_ZIP_CONTENTS = True
MAX_ZIP_MEMBERS = 100

# Extensions whose content can be decoded as plain text
TEXT_EXTENSIONS = {
    ".txt", ".csv", ".log", ".json", ".xml", ".html", ".htm",
    ".md", ".py", ".js", ".ts", ".java", ".cs", ".sql", ".cfg", ".ini",
}


def _human_size(size_bytes: int) -> str:
    if size_bytes >= 1_048_576:
        return f"{size_bytes / 1_048_576:.1f} MB"
    return f"{size_bytes / 1024:.1f} KB"


def _decode_b64(content_b64: str) -> bytes:
    return base64.b64decode(content_b64)


def _extract_pdf_text(data: bytes) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        parts = []
        for page in reader.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                pass
        return "\n".join(parts)
    except Exception:
        return ""


def _extract_docx_text(data: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception:
        return ""


def _extract_xlsx_text(data: bytes) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                parts.extend(str(cell) for cell in row if cell is not None)
        return " ".join(parts)
    except Exception:
        return ""


def _extract_pptx_text(data: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
        return "\n".join(parts)
    except Exception:
        return ""


def _extract_text(att: AttachmentPayload, data: bytes) -> str:
    """Extract text from an attachment based on its extension / content type."""
    ext = os.path.splitext(att.name)[1].lower()

    if ext in TEXT_EXTENSIONS:
        try:
            return data.decode("utf-8", errors="replace")
        except Exception:
            try:
                return data.decode("latin-1", errors="replace")
            except Exception:
                return ""

    if ext == ".pdf":
        return _extract_pdf_text(data)

    if ext == ".docx":
        return _extract_docx_text(data)

    if ext in (".xlsx", ".xlsm"):
        return _extract_xlsx_text(data)

    if ext in (".pptx", ".pptm"):
        return _extract_pptx_text(data)

    return ""


def check(
    attachments: List[AttachmentPayload],
) -> Tuple[List[Violation], List[str]]:
    """
    Run the six-tier attachment inspection pipeline.

    Returns:
      violations : list of Violation objects
      texts      : list of extracted text strings, one per attachment
                   (empty string if not extracted / not applicable)
    """
    violations: List[Violation] = []
    texts: List[str] = []

    for att in attachments:
        ext = os.path.splitext(att.name)[1].lower()
        extracted_text = ""

        # --- Tier 1: Extension block ------------------------------------
        if ext in BLOCKED_EXTENSIONS:
            violations.append(Violation(
                rule_id="ATTACHMENT_BLOCKED_EXT",
                severity="block",
                title="Blocked attachment type",
                detail=(
                    f"'{att.name}' has a blocked file extension ({ext}). "
                    f"Remove this attachment before sending."
                ),
                affected=[att.name],
            ))
            texts.append(extracted_text)
            continue  # Skip further checks for this file

        # --- Tier 2: Size limits ----------------------------------------
        if att.size_bytes > BLOCK_SIZE_MB * 1_048_576:
            violations.append(Violation(
                rule_id="ATTACHMENT_TOO_LARGE",
                severity="block",
                title="Attachment exceeds size limit",
                detail=(
                    f"'{att.name}' is {_human_size(att.size_bytes)}, "
                    f"exceeding the {BLOCK_SIZE_MB} MB limit."
                ),
                affected=[att.name],
            ))
            texts.append(extracted_text)
            continue

        if att.size_bytes > WARN_SIZE_MB * 1_048_576:
            violations.append(Violation(
                rule_id="ATTACHMENT_LARGE",
                severity="warn",
                title="Large attachment",
                detail=(
                    f"'{att.name}' is {_human_size(att.size_bytes)}. "
                    f"Consider whether this file needs to be sent."
                ),
                affected=[att.name],
            ))

        # --- Warn extension (zip, rar, etc.) ----------------------------
        if ext in WARN_EXTENSIONS:
            violations.append(Violation(
                rule_id="ATTACHMENT_WARN_EXT",
                severity="warn",
                title="Compressed archive attachment",
                detail=(
                    f"'{att.name}' is a compressed archive ({ext}). "
                    f"Ensure the contents are appropriate to send."
                ),
                affected=[att.name],
            ))

        # --- Tiers 3-6: Content extraction (requires Base64 bytes) ------
        if att.content_b64:
            try:
                data = _decode_b64(att.content_b64)
            except Exception:
                violations.append(Violation(
                    rule_id="ATTACHMENT_DECODE_ERROR",
                    severity="warn",
                    title="Attachment could not be decoded",
                    detail=f"'{att.name}' could not be decoded for content scanning.",
                    affected=[att.name],
                ))
                texts.append(extracted_text)
                continue

            # Tier 6: ZIP member inspection
            if ext == ".zip" and SCAN_ZIP_CONTENTS:
                try:
                    with zipfile.ZipFile(io.BytesIO(data)) as zf:
                        members = zf.namelist()
                        if len(members) > MAX_ZIP_MEMBERS:
                            violations.append(Violation(
                                rule_id="ATTACHMENT_ZIP_TOO_MANY",
                                severity="warn",
                                title="Archive contains many files",
                                detail=(
                                    f"'{att.name}' contains {len(members)} files "
                                    f"(limit {MAX_ZIP_MEMBERS}). Manual review recommended."
                                ),
                                affected=[att.name],
                            ))
                        else:
                            for member in members:
                                member_ext = os.path.splitext(member)[1].lower()
                                if member_ext in BLOCKED_EXTENSIONS:
                                    violations.append(Violation(
                                        rule_id="ATTACHMENT_ZIP_BLOCKED_MEMBER",
                                        severity="block",
                                        title="Archive contains blocked file type",
                                        detail=(
                                            f"'{att.name}' contains '{member}' "
                                            f"which has a blocked extension ({member_ext})."
                                        ),
                                        affected=[att.name, member],
                                    ))
                except zipfile.BadZipFile:
                    pass

            # Tiers 3-5: Text extraction for content scanning
            extracted_text = _extract_text(att, data)
            if not extracted_text and ext in {".pdf", ".docx", ".xlsx", ".pptx"}:
                violations.append(Violation(
                    rule_id="ATTACHMENT_UNSCANNABLE",
                    severity="warn",
                    title="Attachment content could not be scanned",
                    detail=(
                        f"'{att.name}' could not be read for content scanning "
                        f"(may be password-protected or image-based). "
                        f"Manually verify the content is appropriate to send."
                    ),
                    affected=[att.name],
                ))

        texts.append(extracted_text)

    return violations, texts
